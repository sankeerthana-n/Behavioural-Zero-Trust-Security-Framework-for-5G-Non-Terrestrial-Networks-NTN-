import sys
from pptx import Presentation
import PyPDF2

def extract_pdf(file_path):
    print(f"\n--- PDF: {file_path} ---")
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    print(f"Page {i+1}:")
                    print(text[:500] + ("..." if len(text) > 500 else "")) # Print first 500 chars per page to avoid overflow
    except Exception as e:
        print(f"Error reading PDF: {e}")

def extract_pptx(file_path):
    print(f"\n--- PPTX: {file_path} ---")
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            text = " ".join(text_runs)
            if text:
                print(f"Slide {i+1}: {text[:500]}")
    except Exception as e:
        print(f"Error reading PPTX: {e}")

if __name__ == "__main__":
    extract_pdf(r"d:\SANKEERTHANA FILES\Mtech Project\final year thesis\zero trust 1\phase 1 thesis 1BM23SCN08.pdf")
    extract_pptx(r"d:\SANKEERTHANA FILES\Mtech Project\final year thesis\zero trust 1\Behavioural Zero-Trust Security Framework for 5G Non-Terrestrial Networks.pptx")
